"""
Export Engine — HTML, PDF, PPT
Charts are rendered to PNG using plotly's SVG → cairosvg pipeline
or matplotlib fallback (no Chrome required).
"""
import io
import plotly.io as pio


# ─────────────────────────────────────────────────────────────
# CHART → PNG bytes  (no Chrome required)
# ─────────────────────────────────────────────────────────────

def _fig_to_png(fig, width=600, height=380):
    """
    Convert plotly figure to PNG bytes.
    Uses kaleido v0.2.1 (no Chrome needed).
    Falls back to matplotlib if kaleido unavailable.
    """
    try:
        fig.update_layout(width=width, height=height, margin=dict(t=40,b=30,l=30,r=20))
        return fig.to_image(format="png")
    except Exception:
        pass

    # Matplotlib fallback
    from chart_renderer import fig_to_png
    try:
        return fig_to_png(fig.to_json(), width_px=width, height_px=height)
    except Exception:
        return None


def _get_fig(chart):
    try:
        if "fig_json" in chart:
            return pio.from_json(chart["fig_json"])
        return chart.get("fig")
    except Exception:
        return None


# ─────────────────────────────────────────────────────────────
# HTML EXPORT
# ─────────────────────────────────────────────────────────────

def export_html(report_name, rpt_kpis, rpt_charts, insights):

    kpi_html = ""
    for k, v in list(rpt_kpis.items())[:5]:
        kpi_html += f'''
        <div class="kpi-card">
            <div class="kpi-label">{k}</div>
            <div class="kpi-value">{v}</div>
        </div>'''

    charts_html = ""
    for chart in rpt_charts:
        fig = _get_fig(chart)
        if fig:
            ch = pio.to_html(fig, full_html=False, include_plotlyjs="cdn")
            charts_html += f'<div class="chart-box"><b>{chart["title"]}</b><br>{ch}</div>'

    insights_html = ""
    for ins in insights:
        clean = ins.replace("**","").replace("*","")
        insights_html += f'<div class="insight">{clean}</div>'

    html = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"><title>{report_name}</title>
<style>
  body{{font-family:Arial,sans-serif;padding:30px;background:#f4f6f9;margin:0;}}
  h1{{color:#2c3e50;border-bottom:3px solid #3498db;padding-bottom:10px;}}
  h2{{color:#34495e;margin-top:30px;}}
  .kpi-row{{display:flex;gap:16px;flex-wrap:wrap;margin-bottom:28px;}}
  .kpi-card{{background:white;border-radius:12px;padding:18px 28px;
             box-shadow:0 2px 8px rgba(0,0,0,0.08);text-align:center;min-width:120px;}}
  .kpi-label{{font-size:12px;color:#888;margin-bottom:6px;}}
  .kpi-value{{font-size:24px;font-weight:bold;color:#2c3e50;}}
  .chart-grid{{display:grid;grid-template-columns:1fr 1fr;gap:20px;margin-bottom:28px;}}
  .chart-box{{background:white;border-radius:12px;padding:16px;
              box-shadow:0 2px 8px rgba(0,0,0,0.08);}}
  .insight{{background:#eaf4fb;border-left:4px solid #3498db;
            padding:10px 15px;margin:8px 0;border-radius:4px;font-size:14px;}}
  @media(max-width:700px){{.chart-grid{{grid-template-columns:1fr;}}}}
</style></head><body>
<h1>📊 {report_name}</h1>
<div class="kpi-row">{kpi_html}</div>
<h2>📈 Charts</h2>
<div class="chart-grid">{charts_html}</div>
<h2>💡 Insights</h2>
{insights_html}
</body></html>"""

    return html.encode("utf-8")


# ─────────────────────────────────────────────────────────────
# PDF EXPORT  (reportlab + matplotlib chart rendering)
# ─────────────────────────────────────────────────────────────

def export_pdf(report_name, rpt_kpis, rpt_charts, insights):
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer,
        Table, TableStyle, Image as RLImage, HRFlowable
    )

    buf    = io.BytesIO()
    doc    = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=2*cm, rightMargin=2*cm,
        topMargin=2*cm, bottomMargin=2*cm,
    )
    styles = getSampleStyleSheet()
    story  = []

    # Title
    title_style = ParagraphStyle(
        "ReportTitle", parent=styles["Title"],
        fontSize=22, textColor=colors.HexColor("#2c3e50"),
        spaceAfter=6,
    )
    story.append(Paragraph(report_name, title_style))
    story.append(HRFlowable(width="100%", thickness=2, color=colors.HexColor("#3498db")))
    story.append(Spacer(1, 0.4*cm))

    # KPI table
    kpi_items = list(rpt_kpis.items())[:5]
    if kpi_items:
        kpi_data = [[
            Paragraph(
                f'<b>{v}</b><br/><font size="8" color="#888888">{k}</font>',
                styles["Normal"]
            )
            for k, v in kpi_items
        ]]
        col_w = [(A4[0] - 4*cm) / len(kpi_items)] * len(kpi_items)
        kpi_table = Table(kpi_data, colWidths=col_w)
        kpi_table.setStyle(TableStyle([
            ("BACKGROUND",   (0,0), (-1,-1), colors.HexColor("#f8f9fa")),
            ("BOX",          (0,0), (-1,-1), 0.5, colors.HexColor("#dee2e6")),
            ("INNERGRID",    (0,0), (-1,-1), 0.5, colors.HexColor("#dee2e6")),
            ("ALIGN",        (0,0), (-1,-1), "CENTER"),
            ("VALIGN",       (0,0), (-1,-1), "MIDDLE"),
            ("TOPPADDING",   (0,0), (-1,-1), 10),
            ("BOTTOMPADDING",(0,0), (-1,-1), 10),
        ]))
        story.append(kpi_table)
        story.append(Spacer(1, 0.5*cm))

    # Charts
    h2 = ParagraphStyle("H2", parent=styles["Heading2"],
                        fontSize=14, textColor=colors.HexColor("#34495e"),
                        spaceBefore=10, spaceAfter=6)
    story.append(Paragraph("Charts", h2))

    page_w = A4[0] - 4*cm
    img_w  = (page_w - 0.4*cm) / 2
    img_h  = img_w * 0.6

    for i in range(0, len(rpt_charts), 2):
        pair     = rpt_charts[i:i+2]
        row_imgs = []

        for chart in pair:
            fig = _get_fig(chart)
            if fig:
                png = _fig_to_png(fig, width=560, height=340)
                if png:
                    row_imgs.append(RLImage(io.BytesIO(png), width=img_w, height=img_h))
                else:
                    row_imgs.append(Paragraph(chart["title"], styles["Normal"]))
            else:
                row_imgs.append(Spacer(img_w, img_h))

        while len(row_imgs) < 2:
            row_imgs.append(Spacer(img_w, img_h))

        tbl = Table([row_imgs], colWidths=[img_w, img_w])
        tbl.setStyle(TableStyle([
            ("VALIGN",       (0,0), (-1,-1), "TOP"),
            ("LEFTPADDING",  (0,0), (-1,-1), 3),
            ("RIGHTPADDING", (0,0), (-1,-1), 3),
        ]))
        story.append(tbl)
        story.append(Spacer(1, 0.3*cm))

    # Insights
    story.append(Paragraph("Insights", h2))
    ins_style = ParagraphStyle(
        "Insight", parent=styles["Normal"],
        fontSize=10, leading=14,
        leftIndent=8, spaceAfter=5,
    )
    for ins in insights:
        clean = ins.replace("**","").replace("*","")
        story.append(Paragraph(f"• {clean}", ins_style))

    doc.build(story)
    return buf.getvalue()


# ─────────────────────────────────────────────────────────────
# PPT EXPORT  (python-pptx)
# ─────────────────────────────────────────────────────────────

def export_ppt(report_name, rpt_kpis, rpt_charts, insights):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    DARK  = RGBColor(0x2c, 0x3e, 0x50)
    BLUE  = RGBColor(0x34, 0x98, 0xdb)
    LIGHT = RGBColor(0xf4, 0xf6, 0xf9)
    WHITE = RGBColor(0xff, 0xff, 0xff)
    GRAY  = RGBColor(0x88, 0x88, 0x88)
    blank = prs.slide_layouts[6]

    # ── Slide 1: Title + KPIs + Insights ─────────────────────
    sl = prs.slides.add_slide(blank)
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = LIGHT

    # Title bar
    bar = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.3))
    bar.fill.solid(); bar.fill.fore_color.rgb = DARK
    bar.line.fill.background()

    tf = sl.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.5), Inches(1.0))
    p  = tf.text_frame.paragraphs[0]
    r  = p.add_run()
    r.text = report_name
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = WHITE

    # KPI cards
    kpi_items = list(rpt_kpis.items())[:5]
    n         = len(kpi_items)
    card_w    = Inches(2.2)
    card_h    = Inches(1.2)
    gap       = Inches(0.25)
    total_w   = n * card_w + (n-1) * gap
    start_x   = (Inches(13.33) - total_w) / 2

    for i, (k, v) in enumerate(kpi_items):
        cx = start_x + i * (card_w + gap)
        cy = Inches(1.5)

        card = sl.shapes.add_shape(1, cx, cy, card_w, card_h)
        card.fill.solid(); card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(0xde, 0xe2, 0xe6)
        card.line.width = Pt(1)

        vb = sl.shapes.add_textbox(cx+Inches(0.1), cy+Inches(0.1), card_w-Inches(0.2), Inches(0.65))
        vp = vb.text_frame.paragraphs[0]; vp.alignment = PP_ALIGN.CENTER
        vr = vp.add_run(); vr.text = str(v)
        vr.font.size = Pt(20); vr.font.bold = True; vr.font.color.rgb = DARK

        lb = sl.shapes.add_textbox(cx+Inches(0.1), cy+Inches(0.72), card_w-Inches(0.2), Inches(0.4))
        lp = lb.text_frame.paragraphs[0]; lp.alignment = PP_ALIGN.CENTER
        lr = lp.add_run(); lr.text = k
        lr.font.size = Pt(9); lr.font.color.rgb = GRAY

    # Insights
    ib = sl.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(12.3), Inches(4.3))
    ib.text_frame.word_wrap = True
    for j, ins in enumerate(insights[:7]):
        clean = ins.replace("**","").replace("*","")
        para  = ib.text_frame.paragraphs[0] if j == 0 else ib.text_frame.add_paragraph()
        run   = para.add_run()
        run.text = f"• {clean}"
        run.font.size = Pt(10); run.font.color.rgb = DARK
        para.space_after = Pt(3)

    # ── Slides 2+: 2 charts per slide ────────────────────────
    for i in range(0, len(rpt_charts), 2):
        pair = rpt_charts[i:i+2]
        sl   = prs.slides.add_slide(blank)
        sl.background.fill.solid()
        sl.background.fill.fore_color.rgb = LIGHT

        # Thin header
        bar = sl.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.55))
        bar.fill.solid(); bar.fill.fore_color.rgb = DARK
        bar.line.fill.background()

        positions = [(Inches(0.2), Inches(0.65)), (Inches(6.75), Inches(0.65))]

        for chart, (cx, cy) in zip(pair, positions):
            fig = _get_fig(chart)

            # Chart title
            ttb = sl.shapes.add_textbox(cx, cy, Inches(6.2), Inches(0.4))
            ttp = ttb.text_frame.paragraphs[0]
            ttr = ttp.add_run()
            ttr.text = chart["title"]
            ttr.font.size = Pt(11); ttr.font.bold = True; ttr.font.color.rgb = DARK

            if fig:
                png = _fig_to_png(fig, width=580, height=360)
                if png:
                    sl.shapes.add_picture(
                        io.BytesIO(png),
                        cx, cy + Inches(0.42),
                        Inches(6.2), Inches(3.85)
                    )

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
